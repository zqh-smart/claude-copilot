from io import BytesIO

from docx import Document as DocxDocument
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter

from app.pipeline.feature_pipeline.chunking.service import ChunkingService
from app.pipeline.feature_pipeline.parser.doc_parser import DocDocumentParser
from app.pipeline.feature_pipeline.parser.docx_parser import DocxDocumentParser
from app.pipeline.feature_pipeline.parser.extract_processor import ExtractProcessor
from app.pipeline.feature_pipeline.parser.html_parser import HtmlDocumentParser
from app.pipeline.feature_pipeline.evaluation.service import ParseEvaluationBenchmarkService
from app.pipeline.feature_pipeline.parser.markdown_parser import MarkdownDocumentParser
from app.pipeline.feature_pipeline.parser.pdf_parser import MineruParseArtifacts, PdfDocumentParser, PdfPageProfile
from app.pipeline.feature_pipeline.parser.ppt_parser import PptDocumentParser
from app.pipeline.feature_pipeline.parser.pptx_parser import PptxDocumentParser
from app.pipeline.feature_pipeline.parser.spreadsheet_parser import SpreadsheetDocumentParser
from app.pipeline.feature_pipeline.schema_mapping.service import FinancialSchemaMappingService
from app.pipeline.feature_pipeline.structure_reconstruction.service import StructureReconstructionService
from app.pipeline.feature_pipeline.segmentation.service import SemanticSegmentationService
from app.pipeline.feature_pipeline.table_intelligence.service import TableIntelligenceService
from src.claude_copilot.schemas.document import DocumentMetadata, ParsedDocument, ParsedPageBlock, ParsedSection, ParsedTable


def build_metadata(filename: str, extension: str) -> DocumentMetadata:
    return DocumentMetadata(
        doc_type="annual_report",
        source="test",
        filename=filename,
        extension=extension,
    )


def run_document_ai_postprocessing(document: ParsedDocument) -> ParsedDocument:
    document = SemanticSegmentationService().segment(document)
    document = TableIntelligenceService().enhance(document)
    document = StructureReconstructionService().reconstruct(document)
    document = FinancialSchemaMappingService().map(document)
    return document


def test_markdown_parser_splits_sections_by_headers() -> None:
    parser = MarkdownDocumentParser(split_by_headers=True)
    content = b"# Overview\nRevenue grew.\n## Risk\nLiquidity pressure increased."

    result = parser.parse(
        doc_id="doc-1",
        content=content,
        metadata=build_metadata("report.md", ".md"),
    )

    assert result.metadata.parse_backend == "native-markdown"
    assert result.metadata.parse_route == "native_markdown"
    assert len(result.sections) == 2
    assert result.sections[0].title == "Overview"
    assert "Revenue grew." in result.sections[0].content
    assert result.sections[1].title == "Risk"


def test_docx_parser_extracts_headings_and_tables() -> None:
    doc = DocxDocument()
    doc.add_heading("Management Discussion", level=1)
    doc.add_paragraph("Operating income improved.")
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Metric"
    table.rows[0].cells[1].text = "Value"
    table.rows[1].cells[0].text = "Revenue"
    table.rows[1].cells[1].text = "100"

    buffer = BytesIO()
    doc.save(buffer)

    parser = DocxDocumentParser(extract_tables=True)
    result = parser.parse(
        doc_id="doc-2",
        content=buffer.getvalue(),
        metadata=build_metadata("report.docx", ".docx"),
    )

    assert result.metadata.parse_backend == "native-docx"
    assert result.metadata.parse_route == "native_docx"
    assert result.sections[0].title == "Management Discussion"
    assert "Operating income improved." in result.sections[0].content
    assert len(result.tables) == 1
    assert "Revenue" in (result.tables[0].raw_markdown or "")
    assert result.tables[0].title == "Management Discussion"
    assert result.tables[0].headers == ["Metric", "Value"]
    assert result.tables[0].rows == [["Revenue", "100"]]
    heading_order = next(block.order for block in result.page_blocks if block.block_type == "heading")
    table_order = next(block.order for block in result.page_blocks if block.block_type == "table")
    assert heading_order < table_order


def test_html_parser_extracts_sections_and_tables() -> None:
    parser = HtmlDocumentParser()
    content = b"""
    <html>
      <head><title>ACME Filing</title></head>
      <body>
        <h1>Overview</h1>
        <p>Revenue grew year over year.</p>
        <h2>Highlights</h2>
        <ul><li>Margin improved</li></ul>
        <table>
          <caption>Key Metrics</caption>
          <tr><th>Metric</th><th>Value</th></tr>
          <tr><td>Revenue</td><td>100</td></tr>
        </table>
      </body>
    </html>
    """

    result = parser.parse(
        doc_id="doc-html-1",
        content=content,
        metadata=build_metadata("filing.html", ".html"),
    )

    assert result.metadata.parse_backend == "native-html"
    assert result.metadata.parse_route == "native_html"
    assert len(result.sections) == 2
    assert result.sections[0].title == "Overview"
    assert "Revenue grew year over year." in result.sections[0].content
    assert result.sections[1].title == "Highlights"
    assert "Margin improved" in result.sections[1].content
    assert len(result.tables) == 1
    assert result.tables[0].title == "Key Metrics"
    assert result.tables[0].headers == ["Metric", "Value"]
    assert result.tables[0].rows == [["Revenue", "100"]]


def test_html_parser_binds_context_heading_to_table_without_caption() -> None:
    parser = HtmlDocumentParser()
    content = b"""
    <html>
      <body>
        <div><strong>CONSOLIDATED BALANCE SHEETS</strong></div>
        <table>
          <tr><th>Metric</th><th>2024</th><th>2023</th></tr>
          <tr><td>Total assets</td><td>400</td><td>390</td></tr>
        </table>
      </body>
    </html>
    """

    result = parser.parse(
        doc_id="doc-html-2",
        content=content,
        metadata=build_metadata("balance.html", ".html"),
    )

    assert any(block.block_type == "heading" and "CONSOLIDATED BALANCE SHEETS" in block.text for block in result.page_blocks)
    assert result.tables[0].title == "CONSOLIDATED BALANCE SHEETS"
    assert result.tables[0].page == 1
    assert result.tables[0].metadata["title_source"] in {"context", "page_block_heading"}


def test_spreadsheet_parser_extracts_sheet_sections_and_tables() -> None:
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Income Statement"
    worksheet.append(["Metric", "2024", "2023"])
    worksheet.append(["Revenue", 100, 90])
    worksheet.append([])
    worksheet.append(["Region", "APAC", "EMEA"])
    worksheet.append(["Share", "40%", "35%"])

    buffer = BytesIO()
    workbook.save(buffer)

    parser = SpreadsheetDocumentParser()
    result = parser.parse(
        doc_id="doc-xlsx-1",
        content=buffer.getvalue(),
        metadata=build_metadata("statement.xlsx", ".xlsx"),
    )

    assert result.metadata.parse_backend == "native-openpyxl"
    assert result.metadata.parse_route == "native_xlsx"
    assert result.metadata.page_count == 1
    assert len(result.sections) == 1
    assert result.sections[0].title == "Income Statement"
    assert "Revenue\t100\t90" in result.sections[0].content
    assert len(result.tables) == 2
    assert result.tables[0].headers == ["Metric", "2024", "2023"]
    assert result.tables[0].rows == [["Revenue", "100", "90"]]
    assert result.tables[1].headers == ["Region", "APAC", "EMEA"]


def test_spreadsheet_parser_emits_heading_and_table_page_blocks() -> None:
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Financial Statements"
    worksheet.append(["Consolidated statements of income"])
    worksheet.append([])
    worksheet.append(["Metric", "2024", "2023"])
    worksheet.append(["Revenue", 100, 90])

    buffer = BytesIO()
    workbook.save(buffer)

    result = SpreadsheetDocumentParser().parse(
        doc_id="doc-xlsx-2",
        content=buffer.getvalue(),
        metadata=build_metadata("financials.xlsx", ".xlsx"),
    )

    heading_blocks = [block for block in result.page_blocks if block.block_type == "heading"]
    table_blocks = [block for block in result.page_blocks if block.block_type == "table"]

    assert any(block.text == "Financial Statements" for block in heading_blocks)
    assert any(block.text == "Consolidated statements of income" for block in heading_blocks)
    assert len(table_blocks) == 1
    assert result.tables[0].title == "Consolidated statements of income"
    assert result.tables[0].source_block_id == table_blocks[0].block_id


def test_spreadsheet_parser_uses_xls_route_when_extension_is_xls(monkeypatch) -> None:
    parser = SpreadsheetDocumentParser()
    monkeypatch.setattr(
        parser,
        "_load_xls",
        lambda _content: [
            {
                "name": "Legacy Sheet",
                "index": 1,
                "rows": [["Metric", "Value"], ["Profit", "30"]],
            }
        ],
    )

    result = parser.parse(
        doc_id="doc-xls-1",
        content=b"legacy-xls",
        metadata=build_metadata("statement.xls", ".xls"),
    )

    assert result.metadata.parse_backend == "native-xls"
    assert result.metadata.parse_route == "native_xls"
    assert len(result.tables) == 1
    assert result.tables[0].rows == [["Profit", "30"]]


def test_pptx_parser_extracts_slide_text_notes_and_tables() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Quarterly Review"

    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(5), Inches(1))
    textbox.text_frame.text = "Revenue improved across regions."

    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = "Discuss operating leverage."

    table = slide.shapes.add_table(2, 2, Inches(1), Inches(2.2), Inches(4), Inches(1.2)).table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Revenue"
    table.cell(1, 1).text = "100"

    buffer = BytesIO()
    presentation.save(buffer)

    parser = PptxDocumentParser()
    result = parser.parse(
        doc_id="doc-pptx-1",
        content=buffer.getvalue(),
        metadata=build_metadata("deck.pptx", ".pptx"),
    )

    assert result.metadata.parse_backend == "native-pptx"
    assert result.metadata.parse_route == "native_pptx"
    assert result.metadata.page_count == 1
    assert len(result.sections) == 1
    assert result.sections[0].title == "Quarterly Review"
    assert "Revenue improved across regions." in result.sections[0].content
    assert "Discuss operating leverage." in result.sections[0].content
    assert len(result.tables) == 1
    assert result.tables[0].headers == ["Metric", "Value"]
    assert result.tables[0].rows == [["Revenue", "100"]]
    assert result.tables[0].title == "Quarterly Review"
    assert any(block.block_type == "heading" and block.text == "Quarterly Review" for block in result.page_blocks)
    assert any(block.block_type == "table" for block in result.page_blocks)


def test_doc_parser_uses_converted_docx_bytes(monkeypatch) -> None:
    doc = DocxDocument()
    doc.add_heading("Legacy Report", level=1)
    doc.add_paragraph("Converted from doc.")
    buffer = BytesIO()
    doc.save(buffer)

    parser = DocDocumentParser()
    monkeypatch.setattr(parser, "_convert_doc_to_docx_content", lambda _content: buffer.getvalue())

    result = parser.parse(
        doc_id="doc-doc-1",
        content=b"legacy-doc",
        metadata=build_metadata("legacy.doc", ".doc"),
    )

    assert result.metadata.parse_backend == "win32com-word"
    assert result.metadata.parse_route == "native_doc"
    assert result.sections[0].title == "Legacy Report"
    assert "Converted from doc." in result.sections[0].content


def test_ppt_parser_uses_converted_pptx_bytes(monkeypatch) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Legacy Deck"
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(4), Inches(1))
    textbox.text_frame.text = "Converted from ppt."

    buffer = BytesIO()
    presentation.save(buffer)

    parser = PptDocumentParser()
    monkeypatch.setattr(parser, "_convert_ppt_to_pptx_content", lambda _content: buffer.getvalue())

    result = parser.parse(
        doc_id="doc-ppt-1",
        content=b"legacy-ppt",
        metadata=build_metadata("legacy.ppt", ".ppt"),
    )

    assert result.metadata.parse_backend == "win32com-powerpoint"
    assert result.metadata.parse_route == "native_ppt"
    assert result.sections[0].title == "Legacy Deck"
    assert "Converted from ppt." in result.sections[0].content


def test_extract_processor_routes_extended_document_types() -> None:
    processor = ExtractProcessor()

    assert ".html" in processor._handlers
    assert ".htm" in processor._handlers
    assert ".doc" in processor._handlers
    assert ".docm" in processor._handlers
    assert ".xls" in processor._handlers
    assert ".xlsx" in processor._handlers
    assert ".ppt" in processor._handlers
    assert ".pptx" in processor._handlers


def test_html_document_enters_document_ai_postprocessing_chain() -> None:
    parser = HtmlDocumentParser()
    content = b"""
    <html>
      <body>
        <h1>Consolidated statements of income</h1>
        <p>(in millions)</p>
        <table>
          <tr><th>Metric</th><th>2024</th><th>2023</th></tr>
          <tr><td>Revenue</td><td>100</td><td>90</td></tr>
          <tr><td>Net income</td><td>20</td><td>18</td></tr>
        </table>
      </body>
    </html>
    """

    document = parser.parse(
        doc_id="doc-html-chain-1",
        content=content,
        metadata=build_metadata("income.html", ".html"),
    )
    document = run_document_ai_postprocessing(document)

    semantic_sections = [section for section in document.sections if section.metadata.get("source") == "semantic_segmentation"]
    assert any(section.section_type == "financial_statement" for section in semantic_sections)
    assert document.tables[0].table_type == "income_statement"
    assert document.tables[0].normalized_metrics["revenue"] == {"2024": 100, "2023": 90}
    assert document.financial_schema is not None
    assert len(document.financial_schema.statements) == 1


def test_docx_document_enters_document_ai_postprocessing_chain() -> None:
    doc = DocxDocument()
    doc.add_heading("Consolidated balance sheets", level=1)
    doc.add_paragraph("(in millions)")
    table = doc.add_table(rows=3, cols=3)
    table.rows[0].cells[0].text = "Metric"
    table.rows[0].cells[1].text = "2024"
    table.rows[0].cells[2].text = "2023"
    table.rows[1].cells[0].text = "Total assets"
    table.rows[1].cells[1].text = "400"
    table.rows[1].cells[2].text = "390"
    table.rows[2].cells[0].text = "Total liabilities"
    table.rows[2].cells[1].text = "350"
    table.rows[2].cells[2].text = "340"

    buffer = BytesIO()
    doc.save(buffer)

    document = DocxDocumentParser(extract_tables=True).parse(
        doc_id="doc-docx-chain-1",
        content=buffer.getvalue(),
        metadata=build_metadata("balance.docx", ".docx"),
    )
    document = run_document_ai_postprocessing(document)

    semantic_sections = [section for section in document.sections if section.metadata.get("source") == "semantic_segmentation"]
    assert any(section.section_type == "financial_statement" for section in semantic_sections)
    assert document.tables[0].table_type == "balance_sheet"
    assert document.tables[0].normalized_metrics["total_assets"] == {"2024": 400, "2023": 390}
    assert document.financial_schema is not None
    assert len(document.financial_schema.statements) == 1


def test_spreadsheet_document_enters_document_ai_postprocessing_chain() -> None:
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Statements of income"
    worksheet.append(["Metric", "2024", "2023"])
    worksheet.append(["Revenue", 100, 90])
    worksheet.append(["Net income", 20, 18])

    buffer = BytesIO()
    workbook.save(buffer)

    document = SpreadsheetDocumentParser().parse(
        doc_id="doc-xlsx-chain-1",
        content=buffer.getvalue(),
        metadata=build_metadata("income.xlsx", ".xlsx"),
    )
    document = run_document_ai_postprocessing(document)

    semantic_sections = [section for section in document.sections if section.metadata.get("source") == "semantic_segmentation"]
    assert any(section.section_type == "financial_statement" for section in semantic_sections)
    assert document.tables[0].source_section == "financial_statement"
    assert document.tables[0].table_type == "income_statement"
    assert document.tables[0].normalized_metrics["net_income"] == {"2024": 20, "2023": 18}
    assert document.financial_schema is not None
    assert len(document.financial_schema.metric_facts) >= 4


def test_pptx_document_enters_document_ai_postprocessing_chain() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Consolidated statements of income"
    table = slide.shapes.add_table(3, 3, Inches(1), Inches(1.5), Inches(5), Inches(1.5)).table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "2024"
    table.cell(0, 2).text = "2023"
    table.cell(1, 0).text = "Revenue"
    table.cell(1, 1).text = "100"
    table.cell(1, 2).text = "90"
    table.cell(2, 0).text = "Net income"
    table.cell(2, 1).text = "20"
    table.cell(2, 2).text = "18"

    buffer = BytesIO()
    presentation.save(buffer)

    document = PptxDocumentParser().parse(
        doc_id="doc-pptx-chain-1",
        content=buffer.getvalue(),
        metadata=build_metadata("income.pptx", ".pptx"),
    )
    document = run_document_ai_postprocessing(document)

    semantic_sections = [section for section in document.sections if section.metadata.get("source") == "semantic_segmentation"]
    assert any(section.section_type == "financial_statement" for section in semantic_sections)
    assert document.tables[0].table_type == "income_statement"
    assert document.tables[0].normalized_metrics["revenue"] == {"2024": 100, "2023": 90}
    assert document.financial_schema is not None
    assert len(document.financial_schema.statements) == 1


def test_pdf_parser_forced_native_route_on_blank_pdf() -> None:
    writer = PdfWriter()
    writer.add_blank_page(width=300, height=300)
    buffer = BytesIO()
    writer.write(buffer)

    parser = PdfDocumentParser(backend_priority=["native_pdf"])
    result = parser.parse(
        doc_id="doc-3",
        content=buffer.getvalue(),
        metadata=build_metadata("report.pdf", ".pdf"),
    )

    assert result.metadata.parse_backend == "pypdf"
    assert result.metadata.parse_route == "native_pdf"
    assert result.metadata.page_count == 1
    assert result.metadata.parsed_page_range == (1, 1)
    assert result.metadata.parsed_page_count == 1
    assert len(result.sections) == 1
    assert len(result.page_blocks) == 0
    assert result.quality is not None


def test_pdf_parser_table_route_extracts_tables(monkeypatch) -> None:
    parser = PdfDocumentParser(backend_priority=["table_pdf", "native_pdf"])

    monkeypatch.setattr(
        parser,
        "_build_page_profiles",
        lambda _content: [
            PdfPageProfile(
                page_number=1,
                text="Metric  Value\nRevenue  100\nProfit  30",
                lines=["Metric  Value", "Revenue  100", "Profit  30"],
                table_groups=[["Metric  Value", "Revenue  100", "Profit  30"]],
            )
        ],
    )

    result = parser.parse(
        doc_id="doc-4",
        content=b"mock-pdf-content",
        metadata=build_metadata("table.pdf", ".pdf"),
    )

    assert result.metadata.parse_route == "table_pdf"
    assert result.metadata.parse_backend == "pypdf-table-heuristic"
    assert len(result.tables) == 1
    assert result.tables[0].headers == ["Metric", "Value"]
    assert result.tables[0].rows == [["Revenue", "100"], ["Profit", "30"]]
    assert any(block.block_type == "table" for block in result.page_blocks)


def test_pdf_parser_ocr_route_falls_back_when_backend_missing(monkeypatch) -> None:
    parser = PdfDocumentParser(backend_priority=["ocr_pdf", "native_pdf"])

    monkeypatch.setattr(
        parser,
        "_build_page_profiles",
        lambda _content: [
            PdfPageProfile(page_number=1, text="", lines=[], table_groups=[]),
            PdfPageProfile(page_number=2, text="short", lines=["short"], table_groups=[]),
        ],
    )
    monkeypatch.setattr(parser, "_ocr_extract_page_texts", lambda _content: [])

    result = parser.parse(
        doc_id="doc-5",
        content=b"mock-pdf-content",
        metadata=build_metadata("scan.pdf", ".pdf"),
    )

    assert result.metadata.parse_route == "ocr_pdf"
    assert result.metadata.parse_backend == "pypdf-fallback"
    assert any(issue.code == "ocr_backend_unavailable" for issue in result.issues)
    assert result.quality is not None
    assert result.quality.text_coverage == 0.5


def test_pdf_parser_detects_header_footer_and_captions() -> None:
    parser = PdfDocumentParser(backend_priority=["table_pdf", "native_pdf"])
    page_profiles = [
        PdfPageProfile(
            page_number=1,
            text="\n".join(
                [
                    "ACME Annual Report 2025",
                    "Table 1 Revenue Breakdown",
                    "Metric  Value",
                    "Revenue  100",
                    "Figure 1 Cash Trend",
                    "• Growth remained stable",
                    "Page 1",
                ]
            ),
            lines=[
                "ACME Annual Report 2025",
                "Table 1 Revenue Breakdown",
                "Metric  Value",
                "Revenue  100",
                "Figure 1 Cash Trend",
                "• Growth remained stable",
                "Page 1",
            ],
            table_groups=[["Metric  Value", "Revenue  100"]],
        ),
        PdfPageProfile(
            page_number=2,
            text="\n".join(
                [
                    "ACME Annual Report 2025",
                    "Risk Overview",
                    "Page 2",
                ]
            ),
            lines=[
                "ACME Annual Report 2025",
                "Risk Overview",
                "Page 2",
            ],
            table_groups=[],
        ),
    ]

    blocks = parser._build_page_blocks(doc_id="doc-layout", page_profiles=page_profiles, include_tables=True)
    block_types = {(block.page, block.block_type, block.text) for block in blocks}

    assert (1, "header", "ACME Annual Report 2025") in block_types
    assert (1, "table_caption", "Table 1 Revenue Breakdown") in block_types
    assert (1, "table", "Metric  Value\nRevenue  100") in block_types
    assert (1, "figure_caption", "Figure 1 Cash Trend") in block_types
    assert (1, "list_item", "• Growth remained stable") in block_types
    assert (1, "footer", "Page 1") in block_types
    assert (2, "footer", "Page 2") in block_types


def test_pdf_parser_mineru_route_uses_real_adapter_when_available(monkeypatch) -> None:
    parser = PdfDocumentParser(backend_priority=["mineru_pdf", "native_pdf"])

    monkeypatch.setattr(
        parser,
        "_build_page_profiles",
        lambda _content: [
            PdfPageProfile(page_number=1, text="Native fallback text", lines=["Native fallback text"], table_groups=[])
        ],
    )
    monkeypatch.setattr(parser, "_can_use_mineru", lambda: True)
    monkeypatch.setattr(
        parser,
        "_run_mineru_parse",
        lambda **_kwargs: MineruParseArtifacts(
            markdown_text="# Page 1\nRevenue improved",
            page_blocks=[
                ParsedPageBlock(
                    block_id="doc-7-mineru-page-1-block-1",
                    block_type="heading",
                    text="Revenue Overview",
                    page=1,
                    order=1,
                ),
                ParsedPageBlock(
                    block_id="doc-7-mineru-page-1-block-2",
                    block_type="table",
                    text="Metric | Value\nRevenue | 100",
                    page=1,
                    order=2,
                ),
                ParsedPageBlock(
                    block_id="doc-7-mineru-page-1-block-3",
                    block_type="table_caption",
                    text="Table 1 Revenue Breakdown",
                    page=1,
                    order=3,
                ),
            ],
            tables=[
                ParsedTable(
                    table_id="doc-7-mineru-table-1",
                    table_type="mineru_table",
                    title="Table 1 Revenue Breakdown",
                    raw_markdown="| Metric | Value |\n| --- | --- |\n| Revenue | 100 |",
                    page=1,
                    headers=["Metric", "Value"],
                    rows=[["Revenue", "100"]],
                    source_block_id="doc-7-mineru-page-1-block-2",
                    metadata={"parse_mode": "mineru"},
                )
            ],
        ),
    )

    result = parser.parse(
        doc_id="doc-7",
        content=b"mock-pdf-content",
        metadata=build_metadata("mineru.pdf", ".pdf"),
    )

    assert result.metadata.parse_route == "mineru_pdf"
    assert result.metadata.parse_backend == "mineru"
    assert result.metadata.page_count == 1
    assert result.metadata.parsed_page_range == (1, 1)
    assert result.metadata.parsed_page_count == 1
    assert result.raw_text.startswith("# Page 1")
    assert any(block.block_type == "table_caption" for block in result.page_blocks)
    assert result.tables[0].title == "Table 1 Revenue Breakdown"
    assert not any(issue.code == "mineru_backend_unavailable" for issue in result.issues)


def test_pdf_parser_partial_mineru_range_updates_metadata_and_scope(monkeypatch) -> None:
    parser = PdfDocumentParser(
        backend_priority=["mineru_pdf", "native_pdf"],
        mineru_start_page_id=1,
        mineru_end_page_id=2,
    )

    monkeypatch.setattr(
        parser,
        "_build_page_profiles",
        lambda _content: [
            PdfPageProfile(page_number=1, text="Page 1 text", lines=["Page 1 text"], table_groups=[]),
            PdfPageProfile(page_number=2, text="Page 2 text", lines=["Page 2 text"], table_groups=[]),
            PdfPageProfile(page_number=3, text="Page 3 text", lines=["Page 3 text"], table_groups=[]),
        ],
    )
    monkeypatch.setattr(parser, "_can_use_mineru", lambda: True)
    monkeypatch.setattr(
        parser,
        "_run_mineru_parse",
        lambda **_kwargs: MineruParseArtifacts(markdown_text="# Parsed pages 2-3"),
    )

    result = parser.parse(
        doc_id="doc-7b",
        content=b"mock-pdf-content",
        metadata=build_metadata("mineru-partial.pdf", ".pdf"),
    )

    assert result.metadata.parse_route == "mineru_pdf"
    assert result.metadata.parse_backend == "mineru"
    assert result.metadata.page_count == 3
    assert result.metadata.parsed_page_range == (2, 3)
    assert result.metadata.parsed_page_count == 2
    assert len(result.sections) == 2
    assert [section.page_start for section in result.sections] == [2, 3]
    assert result.quality is not None
    assert result.quality.text_coverage == 1.0


def test_pdf_parser_mineru_tables_prefer_canonical_page_payloads_and_dedupe() -> None:
    parser = PdfDocumentParser(backend_priority=["mineru_pdf", "native_pdf"])
    page_blocks = [
        ParsedPageBlock(
            block_id="doc-8-mineru-page-5-block-1",
            block_type="table",
            text="<table><tr><td>Metric</td><td>Value</td></tr><tr><td>Revenue</td><td>100</td></tr></table>",
            page=5,
            order=1,
        )
    ]
    json_payloads = [
        [
            [],
            [],
            [],
            [],
            [
                {
                    "type": "table",
                    "bbox": [10, 20, 30, 40],
                    "content": {
                        "html": "<table><tr><td>Metric</td><td>Value</td></tr><tr><td>Revenue</td><td>100</td></tr></table>",
                        "table_caption": [{"type": "text", "content": "Table 1 Revenue Breakdown"}],
                    },
                }
            ],
        ],
        {
            "pdf_info": [
                {},
                {},
                {},
                {},
                {
                    "para_blocks": [
                        {
                            "type": "table",
                            "bbox": [10, 20, 30, 40],
                            "blocks": [
                                {
                                    "type": "table_body",
                                    "lines": [
                                        {
                                            "spans": [
                                                {
                                                    "type": "table",
                                                    "html": "<table><tr><td>Metric</td><td>Value</td></tr><tr><td>Revenue</td><td>100</td></tr></table>",
                                                }
                                            ]
                                        }
                                    ],
                                }
                            ],
                        }
                    ]
                },
            ]
        },
    ]

    tables = parser._parse_mineru_tables(doc_id="doc-8", json_payloads=json_payloads, page_blocks=page_blocks)

    assert len(tables) == 1
    assert tables[0].page == 5
    assert tables[0].title == "Table 1 Revenue Breakdown"
    assert tables[0].headers == ["Metric", "Value"]
    assert tables[0].rows == [["Revenue", "100"]]


def test_pdf_parser_mineru_tables_bind_caption_and_footnotes_by_bbox() -> None:
    parser = PdfDocumentParser(backend_priority=["mineru_pdf", "native_pdf"])
    table_1_html = "<table><tr><td>Metric</td><td>2024</td></tr><tr><td>Revenue</td><td>100</td></tr></table>"
    table_2_html = "<table><tr><td>Metric</td><td>2024</td></tr><tr><td>Expense</td><td>55</td></tr></table>"
    page_blocks = [
        ParsedPageBlock(
            block_id="doc-8b-page-5-block-1",
            block_type="table_caption",
            text="Table 1 Revenue Breakdown",
            page=5,
            order=1,
            bbox=(40, 100, 560, 120),
        ),
        ParsedPageBlock(
            block_id="doc-8b-page-5-block-2",
            block_type="table",
            text=table_1_html,
            page=5,
            order=2,
            bbox=(40, 130, 560, 280),
        ),
        ParsedPageBlock(
            block_id="doc-8b-page-5-block-3",
            block_type="footnote",
            text="Revenue excludes legacy assets.",
            page=5,
            order=3,
            bbox=(40, 290, 560, 305),
        ),
        ParsedPageBlock(
            block_id="doc-8b-page-5-block-4",
            block_type="table_caption",
            text="Table 2 Expense Breakdown",
            page=5,
            order=4,
            bbox=(40, 340, 560, 360),
        ),
        ParsedPageBlock(
            block_id="doc-8b-page-5-block-5",
            block_type="table",
            text=table_2_html,
            page=5,
            order=5,
            bbox=(40, 370, 560, 520),
        ),
        ParsedPageBlock(
            block_id="doc-8b-page-5-block-6",
            block_type="footnote",
            text="Expense includes litigation charges.",
            page=5,
            order=6,
            bbox=(40, 530, 560, 545),
        ),
    ]
    json_payloads = [
        [
            [],
            [],
            [],
            [],
            [
                {
                    "type": "table",
                    "bbox": [40, 130, 560, 280],
                    "content": {"html": table_1_html},
                },
                {
                    "type": "table",
                    "bbox": [40, 370, 560, 520],
                    "content": {"html": table_2_html},
                },
            ],
        ]
    ]

    tables = parser._parse_mineru_tables(doc_id="doc-8b", json_payloads=json_payloads, page_blocks=page_blocks)

    assert len(tables) == 2
    assert tables[0].title == "Table 1 Revenue Breakdown"
    assert tables[0].footnotes == ["Revenue excludes legacy assets."]
    assert tables[0].metadata["caption_block_id"] == "doc-8b-page-5-block-1"
    assert tables[0].metadata["footnote_block_ids"] == ["doc-8b-page-5-block-3"]
    assert tables[1].title == "Table 2 Expense Breakdown"
    assert tables[1].footnotes == ["Expense includes litigation charges."]
    assert tables[1].metadata["caption_block_id"] == "doc-8b-page-5-block-4"
    assert tables[1].metadata["footnote_block_ids"] == ["doc-8b-page-5-block-6"]


def test_pdf_parser_mineru_page_blocks_prefer_canonical_page_payloads() -> None:
    parser = PdfDocumentParser(backend_priority=["mineru_pdf", "native_pdf"])
    json_payloads = [
        [
            [
                {
                    "type": "title",
                    "bbox": [10, 20, 30, 40],
                    "content": {"title_content": [{"type": "text", "content": "Risk Factors"}]},
                },
                {
                    "type": "paragraph",
                    "bbox": [10, 50, 30, 80],
                    "content": {"paragraph_content": [{"type": "text", "content": "Liquidity pressure remained elevated."}]},
                },
            ]
        ],
        {
            "pdf_info": [
                {
                    "para_blocks": [
                        {
                            "type": "title",
                            "blocks": [
                                {
                                    "lines": [
                                        {"spans": [{"type": "text", "content": "Risk Factors"}]},
                                    ]
                                }
                            ],
                        }
                    ]
                }
            ]
        },
    ]

    blocks = parser._parse_mineru_page_blocks(doc_id="doc-9", json_payloads=json_payloads)

    assert len(blocks) == 2
    assert blocks[0].page == 1
    assert blocks[0].block_type == "heading"
    assert blocks[0].text == "Risk Factors"
    assert blocks[1].page == 1
    assert blocks[1].block_type == "paragraph"


def test_chunking_service_consumes_page_blocks_and_tables() -> None:
    document = ParsedDocument(
        doc_id="doc-6",
        metadata=build_metadata("report.pdf", ".pdf"),
        sections=[
            ParsedSection(
                section_id="doc-6-page-1",
                title="Page 1",
                content="Revenue overview",
                section_type="pdf_page",
                page_start=1,
                page_end=1,
            )
        ],
        page_blocks=[
            ParsedPageBlock(
                block_id="doc-6-page-1-block-1",
                block_type="heading",
                text="Revenue Overview",
                page=1,
                order=1,
            ),
            ParsedPageBlock(
                block_id="doc-6-page-1-block-2",
                block_type="paragraph",
                text="Revenue grew strongly year over year.",
                page=1,
                order=2,
            ),
        ],
        tables=[
            ParsedTable(
                table_id="doc-6-table-1",
                table_type="pdf_table",
                page=1,
                headers=["Metric", "Value"],
                rows=[["Revenue", "100"]],
                raw_markdown="| Metric | Value |\n| --- | --- |\n| Revenue | 100 |",
            )
        ],
    )

    service = ChunkingService(chunk_size=200, overlap=20)
    segments = service.chunk(document)

    assert len(segments) == 3
    assert {segment.metadata["content_type"] for segment in segments} == {"page_block", "table"}
    assert any("Revenue Overview" in segment.content for segment in segments)
    assert any("Metric" in segment.content for segment in segments)


def test_semantic_segmentation_builds_financial_sections_from_page_blocks() -> None:
    service = SemanticSegmentationService()
    document = ParsedDocument(
        doc_id="doc-semantic-1",
        metadata=build_metadata("report.pdf", ".pdf"),
        page_blocks=[
            ParsedPageBlock(
                block_id="doc-semantic-1-page-10-block-1",
                block_type="heading",
                text="Risk Factors",
                page=10,
                order=1,
            ),
            ParsedPageBlock(
                block_id="doc-semantic-1-page-10-block-2",
                block_type="paragraph",
                text="Liquidity pressure could affect funding costs.",
                page=10,
                order=2,
            ),
            ParsedPageBlock(
                block_id="doc-semantic-1-page-11-block-1",
                block_type="heading",
                text="Consolidated statements of income",
                page=11,
                order=1,
            ),
            ParsedPageBlock(
                block_id="doc-semantic-1-page-11-block-2",
                block_type="table",
                text="Revenue | 100",
                page=11,
                order=2,
            ),
        ],
    )

    result = service.segment(document)
    semantic_sections = [section for section in result.sections if section.metadata.get("source") == "semantic_segmentation"]

    assert len(semantic_sections) == 2
    assert semantic_sections[0].section_type == "risk_section"
    assert semantic_sections[0].page_start == 10
    assert semantic_sections[0].page_end == 10
    assert "Liquidity pressure" in semantic_sections[0].content
    assert semantic_sections[1].section_type == "financial_statement"
    assert "Revenue | 100" in semantic_sections[1].content


def test_chunking_service_emits_semantic_section_segments() -> None:
    document = ParsedDocument(
        doc_id="doc-semantic-2",
        metadata=build_metadata("report.pdf", ".pdf"),
        sections=[
            ParsedSection(
                section_id="doc-semantic-2-semantic-section-1",
                title="Risk Factors",
                content="Risk Factors\n\nLiquidity pressure could affect funding costs.",
                section_type="risk_section",
                page_start=10,
                page_end=12,
                metadata={"source": "semantic_segmentation", "confidence": 0.95},
            )
        ],
    )

    service = ChunkingService(chunk_size=500, overlap=50)
    segments = service.chunk(document)

    assert len(segments) == 1
    assert segments[0].metadata["content_type"] == "semantic_section"
    assert segments[0].metadata["section_type"] == "risk_section"
    assert segments[0].metadata["semantic_confidence"] == 0.95
    assert "Liquidity pressure" in segments[0].content


def test_table_intelligence_classifies_financial_statement_and_normalizes_metrics() -> None:
    service = TableIntelligenceService()
    document = ParsedDocument(
        doc_id="doc-table-1",
        metadata=build_metadata("report.pdf", ".pdf"),
        sections=[
            ParsedSection(
                section_id="doc-table-1-semantic-section-1",
                title="JPMorgan Chase & Co. Consolidated statements of income",
                content="...",
                section_type="financial_statement",
                page_start=5,
                page_end=5,
                metadata={"source": "semantic_segmentation", "confidence": 0.98},
            )
        ],
        tables=[
            ParsedTable(
                table_id="doc-table-1-table-1",
                table_type="mineru_table",
                title="JPMorgan Chase & Co. Consolidated statements of income",
                page=5,
                headers=["Year ended December 31, (in millions)", "2024", "", "2023", "2022"],
                rows=[
                    ["Total net revenue", "177,556", "", "158,104", "128,695"],
                    ["Net income", "$ 58,471", "$", "49,552 $", "37,676"],
                ],
                footnotes=["The Notes to Consolidated Financial Statements are an integral part of these statements."],
                raw_markdown="| Year ended December 31, (in millions) | 2024 |  | 2023 | 2022 |\n| --- | --- | --- | --- | --- |\n| Total net revenue | 177,556 |  | 158,104 | 128,695 |\n| Net income | $ 58,471 | $ | 49,552 $ | 37,676 |",
                metadata={"raw_table_type": "complex_table"},
            )
        ],
    )

    result = service.enhance(document)
    table = result.tables[0]

    assert table.table_type == "income_statement"
    assert table.period_headers == ["2024", "2023", "2022"]
    assert table.unit == "millions"
    assert table.currency == "USD"
    assert table.source_section == "financial_statement"
    assert table.normalized_metrics["revenue"] == {"2024": 177556, "2023": 158104, "2022": 128695}
    assert table.normalized_metrics["net_income"] == {"2024": 58471, "2023": 49552, "2022": 37676}
    assert table.footnotes == ["The Notes to Consolidated Financial Statements are an integral part of these statements."]


def test_table_intelligence_normalizes_extended_primary_statement_metrics() -> None:
    service = TableIntelligenceService()
    document = ParsedDocument(
        doc_id="doc-table-extended-1",
        metadata=build_metadata("extended.pdf", ".pdf"),
        sections=[
            ParsedSection(
                section_id="doc-table-extended-1-section-1",
                title="Consolidated statements of cash flows",
                content="...",
                section_type="financial_statement",
                page_start=9,
                page_end=9,
                metadata={"source": "semantic_segmentation", "confidence": 0.98, "anchor_block_id": "doc-table-extended-1-page-9-block-1"},
            ),
            ParsedSection(
                section_id="doc-table-extended-1-section-2",
                title="Consolidated statements of shareholders' equity",
                content="...",
                section_type="financial_statement",
                page_start=10,
                page_end=10,
                metadata={"source": "semantic_segmentation", "confidence": 0.98, "anchor_block_id": "doc-table-extended-1-page-10-block-1"},
            ),
        ],
        page_blocks=[
            ParsedPageBlock(
                block_id="doc-table-extended-1-page-9-block-1",
                block_type="heading",
                text="Consolidated statements of cash flows",
                page=9,
                order=1,
            ),
            ParsedPageBlock(
                block_id="doc-table-extended-1-page-9-block-2",
                block_type="table",
                text="cash flow table",
                page=9,
                order=2,
            ),
            ParsedPageBlock(
                block_id="doc-table-extended-1-page-10-block-1",
                block_type="heading",
                text="Consolidated statements of shareholders' equity",
                page=10,
                order=1,
            ),
            ParsedPageBlock(
                block_id="doc-table-extended-1-page-10-block-2",
                block_type="table",
                text="equity table",
                page=10,
                order=2,
            ),
        ],
        tables=[
            ParsedTable(
                table_id="doc-table-extended-1-table-1",
                table_type="mineru_table",
                title="Consolidated statements of cash flows",
                page=9,
                headers=["Year ended December 31", "2024", "2023"],
                rows=[
                    ["Net income", "50", "45"],
                    ["Depreciation and amortization", "12", "11"],
                    ["Net cash provided by operating activities", "70", "62"],
                    ["Net cash used in investing activities", "(25)", "(22)"],
                    ["Capital expenditures", "(18)", "(16)"],
                ],
                raw_markdown="| Year ended December 31 | 2024 | 2023 |\n| --- | --- | --- |\n| Net income | 50 | 45 |\n| Depreciation and amortization | 12 | 11 |\n| Net cash provided by operating activities | 70 | 62 |\n| Net cash used in investing activities | (25) | (22) |\n| Capital expenditures | (18) | (16) |",
                source_block_id="doc-table-extended-1-page-9-block-2",
            ),
            ParsedTable(
                table_id="doc-table-extended-1-table-2",
                table_type="mineru_table",
                title="Consolidated statements of shareholders' equity",
                page=10,
                headers=["Year ended December 31", "2024", "2023"],
                rows=[
                    ["Common stock", "5", "5"],
                    ["Additional paid-in capital", "110", "105"],
                    ["Treasury stock", "(30)", "(28)"],
                    ["Retained earnings", "220", "200"],
                    ["Accumulated other comprehensive income", "8", "6"],
                    ["Total equity", "313", "288"],
                ],
                raw_markdown="| Year ended December 31 | 2024 | 2023 |\n| --- | --- | --- |\n| Common stock | 5 | 5 |\n| Additional paid-in capital | 110 | 105 |\n| Treasury stock | (30) | (28) |\n| Retained earnings | 220 | 200 |\n| Accumulated other comprehensive income | 8 | 6 |\n| Total equity | 313 | 288 |",
                source_block_id="doc-table-extended-1-page-10-block-2",
            ),
        ],
    )

    result = service.enhance(document)
    cash_flow_table = result.tables[0]
    equity_table = result.tables[1]

    assert cash_flow_table.table_type == "cash_flow"
    assert cash_flow_table.normalized_metrics["net_cash_from_operating_activities"] == {"2024": 70, "2023": 62}
    assert cash_flow_table.normalized_metrics["capital_expenditures"] == {"2024": -18, "2023": -16}
    assert cash_flow_table.normalized_metrics["depreciation_and_amortization"] == {"2024": 12, "2023": 11}

    assert equity_table.table_type == "equity"
    assert equity_table.normalized_metrics["additional_paid_in_capital"] == {"2024": 110, "2023": 105}
    assert equity_table.normalized_metrics["treasury_stock"] == {"2024": -30, "2023": -28}
    assert equity_table.normalized_metrics["total_equity"] == {"2024": 313, "2023": 288}


def test_table_intelligence_merges_cross_page_tables() -> None:
    service = TableIntelligenceService()
    document = ParsedDocument(
        doc_id="doc-table-2",
        metadata=build_metadata("report.pdf", ".pdf"),
        sections=[
            ParsedSection(
                section_id="doc-table-2-semantic-section-1",
                title="JPMorgan Chase & Co. Consolidated statements of income",
                content="...",
                section_type="financial_statement",
                page_start=5,
                page_end=6,
                metadata={"source": "semantic_segmentation", "confidence": 0.98},
            )
        ],
        tables=[
            ParsedTable(
                table_id="doc-table-2-table-1",
                table_type="mineru_table",
                title="JPMorgan Chase & Co. Consolidated statements of income",
                page=5,
                headers=["Year ended December 31, (in millions)", "2024", "2023", "2022"],
                rows=[["Total net revenue", "177,556", "158,104", "128,695"]],
                raw_markdown="| Year ended December 31, (in millions) | 2024 | 2023 | 2022 |\n| --- | --- | --- | --- |\n| Total net revenue | 177,556 | 158,104 | 128,695 |",
            ),
            ParsedTable(
                table_id="doc-table-2-table-2",
                table_type="mineru_table",
                title=None,
                page=6,
                headers=["Year ended December 31, (in millions)", "2024", "2023", "2022"],
                rows=[["Net income", "58,471", "49,552", "37,676"]],
                footnotes=["The Notes to Consolidated Financial Statements are an integral part of these statements."],
                raw_markdown="| Year ended December 31, (in millions) | 2024 | 2023 | 2022 |\n| --- | --- | --- | --- |\n| Net income | 58,471 | 49,552 | 37,676 |",
            ),
        ],
    )

    result = service.enhance(document)

    assert len(result.tables) == 1
    table = result.tables[0]
    assert table.table_type == "income_statement"
    assert table.metadata["page_range"] == [5, 6]
    assert len(table.rows) == 2
    assert table.footnotes == ["The Notes to Consolidated Financial Statements are an integral part of these statements."]
    assert table.normalized_metrics["revenue"]["2024"] == 177556
    assert table.normalized_metrics["net_income"]["2022"] == 37676


def test_table_intelligence_semanticizes_notes_table_rows() -> None:
    service = TableIntelligenceService()
    document = ParsedDocument(
        doc_id="doc-table-3",
        metadata=build_metadata("notes.pdf", ".pdf"),
        sections=[
            ParsedSection(
                section_id="doc-table-3-semantic-section-1",
                title="Note 12 – Allowance for credit losses",
                content="...",
                section_type="financial_note",
                page_start=42,
                page_end=42,
                metadata={"source": "semantic_segmentation", "confidence": 0.97},
            )
        ],
        tables=[
            ParsedTable(
                table_id="doc-table-3-table-1",
                table_type="mineru_table",
                title="Note 12 – Allowance for credit losses",
                page=42,
                headers=["Allowance for credit losses", "Class", "2024", "2023"],
                rows=[
                    ["Ending balance", "Consumer", "1,250", "980"],
                    ["Net charge-offs", "Consumer", "(320)", "(280)"],
                ],
                raw_markdown="| Allowance for credit losses | Class | 2024 | 2023 |\n| --- | --- | --- | --- |\n| Ending balance | Consumer | 1,250 | 980 |\n| Net charge-offs | Consumer | (320) | (280) |",
                metadata={"raw_table_type": "complex_table"},
            )
        ],
    )

    result = service.enhance(document)
    table = result.tables[0]

    assert table.table_type == "notes_table"
    assert table.note_number == "Note 12"
    assert table.note_title == "Allowance for credit losses"
    assert table.note_category == "credit_losses"
    assert table.metadata["note_dimension_headers"] == ["Class"]
    assert len(table.semantic_rows) == 2
    assert table.semantic_rows[0]["dimensions"] == {"Class": "Consumer"}
    assert table.semantic_rows[0]["period_values"] == {"2024": 1250, "2023": 980}
    assert "chargeoff" in table.semantic_rows[1]["tags"]
    assert table.semantic_rows[1]["period_values"] == {"2024": -320, "2023": -280}


def test_financial_schema_mapping_aggregates_statements_notes_and_sections() -> None:
    service = FinancialSchemaMappingService()
    document = ParsedDocument(
        doc_id="doc-schema-1",
        metadata=DocumentMetadata(
            doc_type="annual_report",
            source="test",
            filename="report.pdf",
            extension=".pdf",
            company="JPMorgan Chase & Co.",
            year=2024,
        ),
        sections=[
            ParsedSection(
                section_id="doc-schema-1-semantic-section-1",
                title="Consolidated statements of income",
                content="Income statement evidence",
                section_type="financial_statement",
                page_start=5,
                page_end=6,
                metadata={"source": "semantic_segmentation", "confidence": 0.98, "anchor_block_id": "block-1"},
            ),
            ParsedSection(
                section_id="doc-schema-1-semantic-section-2",
                title="Note 12 – Allowance for credit losses",
                content="Allowance note evidence",
                section_type="financial_note",
                page_start=42,
                page_end=43,
                metadata={"source": "semantic_segmentation", "confidence": 0.97, "anchor_block_id": "block-2"},
            ),
        ],
        tables=[
            ParsedTable(
                table_id="doc-schema-1-table-1",
                table_type="income_statement",
                title="Consolidated statements of income",
                page=5,
                period_headers=["2024", "2023"],
                unit="millions",
                currency="USD",
                normalized_metrics={
                    "revenue": {"2024": 177556, "2023": 158104},
                    "net_income": {"2024": 58471, "2023": 49552},
                },
                footnotes=["The Notes are an integral part of these statements."],
                source_section="financial_statement",
                metadata={
                    "page_range": [5, 6],
                    "source_section_id": "doc-schema-1-semantic-section-1",
                    "parse_mode": "mineru",
                },
            ),
            ParsedTable(
                table_id="doc-schema-1-table-2",
                table_type="notes_table",
                title="Note 12 – Allowance for credit losses",
                page=42,
                note_number="Note 12",
                note_title="Allowance for credit losses",
                note_category="credit_losses",
                period_headers=["2024", "2023"],
                semantic_rows=[
                    {
                        "label": "Ending balance",
                        "label_normalized": "ending balance",
                        "row_type": "metric",
                        "dimensions": {"Class": "Consumer"},
                        "period_values": {"2024": 1250, "2023": 980},
                        "tags": ["balance"],
                    }
                ],
                footnotes=["Allowance excludes purchased credit deteriorated assets."],
                source_section="financial_note",
                metadata={
                    "page_range": [42, 43],
                    "source_section_id": "doc-schema-1-semantic-section-2",
                    "note_dimension_headers": ["Class"],
                    "parse_mode": "mineru",
                },
            ),
        ],
    )

    result = service.map(document)

    assert result.financial_schema is not None
    schema = result.financial_schema
    assert schema.company == "JPMorgan Chase & Co."
    assert schema.year == 2024
    assert schema.reporting_periods == ["2024", "2023"]
    assert len(schema.statements) == 1
    assert schema.statements[0].metrics["revenue"]["2024"] == 177556
    assert schema.statements[0].page_range == (5, 6)
    assert len(schema.notes) == 1
    assert schema.notes[0].note_number == "Note 12"
    assert schema.notes[0].dimension_headers == ["Class"]
    assert schema.notes[0].semantic_rows[0]["period_values"]["2023"] == 980
    assert schema.notes[0].note_facts[0].fact_key == "ending_balance"
    assert len(schema.semantic_sections) == 2
    assert schema.semantic_sections[1].section_type == "financial_note"
    assert len(schema.metric_facts) == 4
    assert len(schema.note_facts) == 1
    assert schema.note_facts[0].fact_key == "ending_balance"
    assert schema.metrics_index["net_income"]["2023"] == 49552
    assert schema.metadata["statement_count"] == 1
    assert schema.metadata["note_count"] == 1
    assert schema.metadata["note_fact_count"] == 1


def test_financial_schema_mapping_builds_stable_note_fact_keys() -> None:
    service = FinancialSchemaMappingService()
    document = ParsedDocument(
        doc_id="doc-schema-2",
        metadata=build_metadata("notes.pdf", ".pdf"),
        tables=[
            ParsedTable(
                table_id="doc-schema-2-table-1",
                table_type="notes_table",
                page=20,
                note_number="Note 12",
                note_title="Allowance for credit losses",
                note_category="credit_losses",
                semantic_rows=[
                    {
                        "label": "Ending balance",
                        "label_normalized": "ending balance",
                        "row_type": "metric",
                        "dimensions": {"Class": "Consumer"},
                        "period_values": {"2024": 1250},
                        "tags": ["balance"],
                    },
                    {
                        "label": "Net charge-offs",
                        "label_normalized": "net charge-offs",
                        "row_type": "metric",
                        "dimensions": {"Class": "Consumer"},
                        "period_values": {"2024": -320},
                        "tags": ["chargeoff"],
                    },
                ],
                source_section="financial_note",
                metadata={"page_range": [20, 20], "parse_mode": "mineru"},
            )
        ],
    )

    result = service.map(document)
    schema = result.financial_schema

    assert schema is not None
    fact_keys = [fact.fact_key for fact in schema.note_facts]
    assert fact_keys == ["ending_balance", "net_charge_offs"]
    assert schema.notes[0].note_facts[1].period_values["2024"] == -320


def test_parse_evaluation_benchmark_reports_schema_coverage() -> None:
    mapping_service = FinancialSchemaMappingService()
    benchmark_service = ParseEvaluationBenchmarkService()
    document = ParsedDocument(
        doc_id="doc-benchmark-1",
        metadata=DocumentMetadata(
            doc_type="annual_report",
            source="test",
            filename="report.pdf",
            extension=".pdf",
            parse_route="mineru_pdf",
            parse_backend="mineru",
            page_count=21,
            parsed_page_range=(1, 21),
            parsed_page_count=21,
        ),
        sections=[
            ParsedSection(
                section_id="doc-benchmark-1-section-1",
                title="Consolidated statements of income",
                content="Income section",
                section_type="financial_statement",
                page_start=5,
                page_end=6,
                metadata={"source": "semantic_segmentation", "confidence": 0.98},
            ),
            ParsedSection(
                section_id="doc-benchmark-1-section-2",
                title="Note 12 – Allowance for credit losses",
                content="Note section",
                section_type="financial_note",
                page_start=10,
                page_end=12,
                metadata={"source": "semantic_segmentation", "confidence": 0.97},
            ),
        ],
        tables=[
            ParsedTable(
                table_id="doc-benchmark-1-table-1",
                table_type="income_statement",
                title="Consolidated statements of income",
                page=5,
                period_headers=["2024", "2023"],
                normalized_metrics={"revenue": {"2024": 177556, "2023": 158104}},
                source_section="financial_statement",
                metadata={"page_range": [5, 6], "parse_mode": "mineru"},
            ),
            ParsedTable(
                table_id="doc-benchmark-1-table-2",
                table_type="notes_table",
                title="Note 12 – Allowance for credit losses",
                page=10,
                note_number="Note 12",
                note_title="Allowance for credit losses",
                note_category="credit_losses",
                semantic_rows=[
                    {
                        "label": "Ending balance",
                        "label_normalized": "ending balance",
                        "row_type": "metric",
                        "dimensions": {"Class": "Consumer"},
                        "period_values": {"2024": 1250},
                        "tags": ["balance"],
                    }
                ],
                source_section="financial_note",
                metadata={"page_range": [10, 12], "parse_mode": "mineru"},
            ),
        ],
    )
    document = mapping_service.map(document)

    report = benchmark_service.evaluate(
        document,
        expected={
            "parse_route": "mineru_pdf",
            "parse_backend": "mineru",
            "min_statements": 1,
            "min_notes": 1,
            "min_metric_facts": 2,
            "min_note_facts": 1,
            "required_statement_types": ["income_statement", "notes_table"],
            "required_semantic_section_types": ["financial_statement", "financial_note"],
            "min_statement_page_range_ratio": 1.0,
            "min_note_domain_facts_ratio": 1.0,
        },
    )

    assert report["counts"]["statements"] == 1
    assert report["counts"]["note_facts"] == 1
    assert report["coverage"]["notes"]["notes_with_domain_facts_ratio"] == 1.0
    assert report["coverage"]["provenance"]["statement_page_range_ratio"] == 1.0
    assert any(check["name"] == "parse_backend" and check["passed"] for check in report["checks"])
    assert report["failures"] == []


def test_table_intelligence_extracts_metrics_from_compact_statement_rows() -> None:
    service = TableIntelligenceService()
    document = ParsedDocument(
        doc_id="doc-table-4",
        metadata=build_metadata("report.pdf", ".pdf"),
        sections=[
            ParsedSection(
                section_id="doc-table-4-semantic-section-1",
                title="Consolidated balance sheets",
                content="...",
                section_type="financial_statement",
                page_start=7,
                page_end=7,
                metadata={"source": "semantic_segmentation", "confidence": 0.98},
            )
        ],
        tables=[
            ParsedTable(
                table_id="doc-table-4-table-1",
                table_type="mineru_table",
                title="Consolidated balance sheets",
                page=7,
                headers=["December 31, (in millions) 2024 2023"],
                rows=[
                    ["Total assets 4,002,000 3,875,000"],
                    ["Total liabilities 3,610,000 3,501,000"],
                ],
                raw_markdown="| December 31, (in millions) 2024 2023 |\n| --- |\n| Total assets 4,002,000 3,875,000 |\n| Total liabilities 3,610,000 3,501,000 |",
            )
        ],
    )

    result = service.enhance(document)
    table = result.tables[0]

    assert table.table_type == "balance_sheet"
    assert table.normalized_metrics["total_assets"] == {"2024": 4002000, "2023": 3875000}
    assert table.normalized_metrics["total_liabilities"] == {"2024": 3610000, "2023": 3501000}


def test_table_intelligence_reads_periods_from_secondary_header_not_maturity_rows() -> None:
    service = TableIntelligenceService()
    document = ParsedDocument(
        doc_id="doc-periods",
        metadata=build_metadata("report.html", ".html"),
        tables=[
            ParsedTable(
                table_id="statement",
                title="Consolidated statements of income",
                headers=["", "Years ended"],
                rows=[
                    ["", "December 31, 2024", "December 31, 2023"],
                    ["Revenue", "100", "90"],
                ],
                raw_markdown="",
                page=1,
            ),
            ParsedTable(
                table_id="maturities",
                title="Note 8 - Lease maturities",
                headers=["", "Operating leases", "Finance leases"],
                rows=[
                    ["2024", "10", "2"],
                    ["2025", "11", "1"],
                ],
                raw_markdown="",
                page=1,
            ),
        ],
    )

    result = service.enhance(document)

    assert result.tables[0].period_headers == ["December 31, 2024", "December 31, 2023"]
    assert result.tables[1].period_headers == []


def test_structure_reconstruction_attaches_note_sections_and_reclassifies_note_tables() -> None:
    service = StructureReconstructionService()
    document = ParsedDocument(
        doc_id="doc-structure-1",
        metadata=build_metadata("report.pdf", ".pdf"),
        sections=[
            ParsedSection(
                section_id="doc-structure-1-note-1",
                title="Note 2 - Fair value measurement",
                content="...",
                section_type="financial_note",
                page_start=14,
                page_end=14,
                metadata={"source": "semantic_segmentation", "confidence": 0.94},
            )
        ],
        tables=[
            ParsedTable(
                table_id="doc-structure-1-table-1",
                table_type="balance_sheet",
                title="Assets and liabilities measured at fair value on a recurring basis",
                page=19,
                metadata={"page_range": [19, 19]},
            ),
            ParsedTable(
                table_id="doc-structure-1-table-2",
                table_type="notes_table",
                title="Notes to consolidated financial statements",
                page=13,
                metadata={"page_range": [13, 13]},
            ),
        ],
    )

    result = service.reconstruct(document)

    assert result.tables[0].source_section == "financial_note"
    assert result.tables[0].table_type == "notes_table"
    assert result.tables[0].note_number == "Note 2"
    assert result.tables[0].note_title == "Fair value measurement"
    assert result.tables[1].title == "Note 2 - Fair value measurement"
    assert result.tables[1].note_number == "Note 2"


def test_financial_schema_mapping_consolidates_statement_fragments() -> None:
    service = FinancialSchemaMappingService()
    document = ParsedDocument(
        doc_id="doc-schema-3",
        metadata=DocumentMetadata(doc_type="annual_report", source="test", filename="report.pdf", extension=".pdf"),
        tables=[
            ParsedTable(
                table_id="doc-schema-3-table-1",
                table_type="income_statement",
                title="Consolidated statements of income",
                page=5,
                period_headers=["2024", "2023"],
                normalized_metrics={"revenue": {"2024": 100, "2023": 90}},
                source_section="financial_statement",
                metadata={"page_range": [5, 5], "parse_mode": "mineru"},
            ),
            ParsedTable(
                table_id="doc-schema-3-table-2",
                table_type="income_statement",
                title="Consolidated statements of income",
                page=7,
                period_headers=["2024", "2023"],
                normalized_metrics={"net_income": {"2024": 20, "2023": 15}},
                source_section="financial_statement",
                metadata={"page_range": [7, 7], "parse_mode": "mineru"},
            ),
        ],
    )

    result = service.map(document)
    schema = result.financial_schema

    assert schema is not None
    assert len(schema.statements) == 1
    assert schema.statements[0].page_range == (5, 7)
    assert schema.statements[0].metrics["revenue"]["2024"] == 100
    assert schema.statements[0].metrics["net_income"]["2023"] == 15
    assert len(schema.metric_facts) == 4
