from __future__ import annotations

from io import BytesIO

from pptx import Presentation

from app.pipeline.feature_pipeline.parser.helpers import with_parse_metadata
from src.claude_copilot.schemas.document import (
    DocumentMetadata,
    ParsedDocument,
    ParsedPageBlock,
    ParsedSection,
    ParsedTable,
)


class PptxDocumentParser:
    def parse(self, *, doc_id: str, content: bytes, metadata: DocumentMetadata) -> ParsedDocument:
        presentation = Presentation(BytesIO(content))
        sections = self._build_sections(doc_id, presentation)
        tables = self._build_tables(doc_id, presentation)
        page_blocks = self._build_page_blocks(doc_id, presentation, tables)

        raw_parts = [section.content for section in sections if section.content.strip()]
        raw_parts.extend(table.raw_markdown or "" for table in tables)
        raw_text = "\n\n".join(part for part in raw_parts if part.strip())

        extension = (metadata.extension or "").lower()
        parse_route = "native_ppt" if extension == ".ppt" else "native_pptx"
        parse_backend = "win32com-powerpoint" if extension == ".ppt" else "native-pptx"

        return ParsedDocument(
            doc_id=doc_id,
            metadata=with_parse_metadata(
                metadata,
                parse_backend=parse_backend,
                parse_route=parse_route,
                page_count=len(presentation.slides),
                parsed_page_range=(1, len(presentation.slides)) if presentation.slides else None,
                parsed_page_count=len(presentation.slides),
            ),
            raw_text=raw_text,
            sections=sections,
            page_blocks=page_blocks,
            tables=tables,
        )

    def _build_sections(self, doc_id: str, presentation: Presentation) -> list[ParsedSection]:
        sections: list[ParsedSection] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            title = None
            if slide.shapes.title and slide.shapes.title.text.strip():
                title = slide.shapes.title.text.strip()

            body_lines: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = shape.text.strip()
                    if text:
                        body_lines.append(text)

            try:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                notes_text = ""
            if notes_text:
                body_lines.append(f"Speaker Notes: {notes_text}")

            if not title:
                title = body_lines[0].splitlines()[0] if body_lines else f"Slide {slide_index}"

            sections.append(
                ParsedSection(
                    section_id=f"{doc_id}-slide-{slide_index}",
                    title=title,
                    content="\n".join(body_lines).strip(),
                    section_type="presentation_slide",
                    page_start=slide_index,
                    page_end=slide_index,
                )
            )

        return sections

    def _build_tables(self, doc_id: str, presentation: Presentation) -> list[ParsedTable]:
        tables: list[ParsedTable] = []
        table_index = 0
        for slide_index, slide in enumerate(presentation.slides, start=1):
            for shape in slide.shapes:
                if not getattr(shape, "has_table", False):
                    continue
                rows = []
                for row in shape.table.rows:
                    rows.append([cell.text.strip().replace("\n", " ") for cell in row.cells])
                if not rows:
                    continue
                table_index += 1
                tables.append(
                    ParsedTable(
                        table_id=f"{doc_id}-pptx-table-{table_index}",
                        table_type="pptx_table",
                        title=slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.text.strip() else f"Slide {slide_index} Table {table_index}",
                        raw_markdown=self._table_to_markdown(rows),
                        page=slide_index,
                        headers=rows[0],
                        rows=rows[1:] if len(rows) > 1 else [],
                        metadata={
                            "row_count": len(rows),
                            "column_count": max((len(row) for row in rows), default=0),
                            "slide_index": slide_index,
                        },
                    )
                )
        return tables

    def _build_page_blocks(
        self,
        doc_id: str,
        presentation: Presentation,
        tables: list[ParsedTable],
    ) -> list[ParsedPageBlock]:
        blocks: list[ParsedPageBlock] = []
        table_lookup = {(table.page, table.title, table.raw_markdown): table for table in tables}

        for slide_index, slide in enumerate(presentation.slides, start=1):
            order = 0
            slide_title = slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.text.strip() else None

            if slide_title:
                order += 1
                blocks.append(
                    ParsedPageBlock(
                        block_id=f"{doc_id}-slide-{slide_index}-block-{order}",
                        block_type="heading",
                        text=slide_title,
                        page=slide_index,
                        order=order,
                    )
                )

            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = shape.text.strip()
                    if not text or text == slide_title:
                        continue
                    order += 1
                    blocks.append(
                        ParsedPageBlock(
                            block_id=f"{doc_id}-slide-{slide_index}-block-{order}",
                            block_type="paragraph",
                            text=text,
                            page=slide_index,
                            order=order,
                        )
                    )
                    continue

                if not getattr(shape, "has_table", False):
                    continue

                rows = [[cell.text.strip().replace("\n", " ") for cell in row.cells] for row in shape.table.rows]
                table_text = self._table_to_markdown(rows)
                order += 1
                block_id = f"{doc_id}-slide-{slide_index}-block-{order}"
                blocks.append(
                    ParsedPageBlock(
                        block_id=block_id,
                        block_type="table",
                        text=table_text,
                        page=slide_index,
                        order=order,
                        metadata={"slide_title": slide_title},
                    )
                )
                matched_table = table_lookup.get((slide_index, slide_title or f"Slide {slide_index} Table {len(table_lookup)}", table_text))
                if matched_table is None:
                    matched_table = next(
                        (
                            table
                            for table in tables
                            if table.page == slide_index and (table.raw_markdown or "") == table_text and not table.source_block_id
                        ),
                        None,
                    )
                if matched_table is not None:
                    matched_table.source_block_id = block_id

            try:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                notes_text = ""
            if notes_text:
                order += 1
                blocks.append(
                    ParsedPageBlock(
                        block_id=f"{doc_id}-slide-{slide_index}-block-{order}",
                        block_type="paragraph",
                        text=f"Speaker Notes: {notes_text}",
                        page=slide_index,
                        order=order,
                        metadata={"source": "speaker_notes"},
                    )
                )

        return blocks

    @staticmethod
    def _table_to_markdown(rows: list[list[str]]) -> str:
        if not rows:
            return ""
        width = max(len(row) for row in rows)
        normalized = [row + [""] * (width - len(row)) for row in rows]
        header = normalized[0]
        lines = [
            "| " + " | ".join(header) + " |",
            "| " + " | ".join(["---"] * width) + " |",
        ]
        for row in normalized[1:]:
            lines.append("| " + " | ".join(row) + " |")
        return "\n".join(lines)
